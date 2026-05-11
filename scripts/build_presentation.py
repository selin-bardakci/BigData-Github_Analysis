#!/usr/bin/env python3
"""
build_presentation.py — Generate a midterm presentation PPTX from the
midterm report content. Photo placeholders are inserted as text boxes
(e.g. "[Image: outputs/d1_timeseries.png]") so the user can replace
them with real images in PowerPoint/Keynote.

Run from project root:
    .venv/bin/python scripts/build_presentation.py
Output:
    report/GitHub_Tech_Trends_Midterm.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

OUT = Path(__file__).resolve().parent.parent / "report" / "GitHub_Tech_Trends_Midterm.pptx"
OUT.parent.mkdir(exist_ok=True)

# ── Theme ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0B, 0x0E, 0x14)
SURFACE  = RGBColor(0x15, 0x1B, 0x27)
INK      = RGBColor(0xE4, 0xE8, 0xF1)
SUB      = RGBColor(0x8B, 0x95, 0xA8)
ACCENT   = RGBColor(0xA7, 0x8B, 0xFA)  # purple
ACCENT2  = RGBColor(0x5E, 0xEA, 0xD4)  # mint
ACCENT3  = RGBColor(0xF4, 0x72, 0xB6)  # pink
ACCENT4  = RGBColor(0xFB, 0xBF, 0x24)  # amber
GOOD     = RGBColor(0x34, 0xD3, 0x99)
BAD      = RGBColor(0xF8, 0x71, 0x71)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────────────────
def add_blank():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.shadow.inherit = False
    return slide

def add_text(slide, text, *, left, top, width, height,
             size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
             font="Inter", italic=False, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if isinstance(text, str):
        lines = [text]
    else:
        lines = list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

def add_bullets(slide, items, *, left, top, width, height,
                size=16, color=INK, bullet_color=ACCENT,
                line_spacing=1.25, indent_size=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet character in accent
        b = p.add_run()
        b.text = "▸ "
        b.font.name = "Inter"
        b.font.size = Pt(size)
        b.font.color.rgb = bullet_color
        b.font.bold = True
        # body text
        r = p.add_run()
        r.text = item
        r.font.name = "Inter"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb

def add_rect(slide, left, top, width, height, fill=SURFACE,
             line=None, corner=False):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh

def header(slide, title, kicker=None):
    add_rect(slide, Inches(0.6), Inches(0.45), Inches(0.08), Inches(0.5),
             fill=ACCENT)
    if kicker:
        add_text(slide, kicker, left=Inches(0.85), top=Inches(0.42),
                 width=Inches(10), height=Inches(0.3),
                 size=11, color=ACCENT, bold=True)
    add_text(slide, title, left=Inches(0.85), top=Inches(0.72),
             width=Inches(11.6), height=Inches(0.7),
             size=28, bold=True, color=INK)
    # underline
    add_rect(slide, Inches(0.85), Inches(1.45),
             Inches(12), Inches(0.015),
             fill=RGBColor(0x1F, 0x27, 0x33))

def footer(slide, n, total):
    add_text(slide, "GitHub Tech Trends · Midterm",
             left=Inches(0.85), top=Inches(7.0),
             width=Inches(8), height=Inches(0.3),
             size=10, color=SUB)
    add_text(slide, f"{n} / {total}",
             left=Inches(12.0), top=Inches(7.0),
             width=Inches(1.0), height=Inches(0.3),
             size=10, color=SUB, align=PP_ALIGN.RIGHT)

def img_placeholder(slide, left, top, width, height, label):
    """Dashed rectangle with the expected image filename for the user."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid(); sh.fill.fore_color.rgb = SURFACE
    sh.line.color.rgb = ACCENT
    sh.line.width = Pt(1.0)
    sh.line.dash_style = 7  # dashed
    sh.shadow.inherit = False
    tb = slide.shapes.add_textbox(left, top + Emu(int(height/2) - 200000),
                                  width, Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "🖼  Image placeholder"
    r.font.name = "Inter"
    r.font.size = Pt(13)
    r.font.color.rgb = SUB
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.name = "JetBrains Mono"
    r2.font.size = Pt(11)
    r2.font.color.rgb = ACCENT2

# ───────────────────────────────────────────────────────────────────────────
# Slides
# ───────────────────────────────────────────────────────────────────────────
SLIDES = []  # collect funcs to build, so we can know total count

def slide_title():
    s = add_blank()
    # decorative bar
    add_rect(s, 0, 0, prs.slide_width, Inches(0.06), fill=ACCENT)
    add_rect(s, 0, Inches(0.06), prs.slide_width, Inches(0.04), fill=ACCENT2)

    add_text(s, "MIDTERM PRESENTATION",
             left=Inches(0.85), top=Inches(0.9),
             width=Inches(8), height=Inches(0.4),
             size=12, bold=True, color=ACCENT)
    add_text(s, "GitHub Technology Trend Analysis",
             left=Inches(0.85), top=Inches(1.3),
             width=Inches(12), height=Inches(1.4),
             size=40, bold=True, color=INK)
    add_text(s, "A Large-Scale Big Data Pipeline for Developer Ecosystem Intelligence",
             left=Inches(0.85), top=Inches(2.5),
             width=Inches(12), height=Inches(0.8),
             size=20, color=ACCENT2, italic=True)

    # author block
    y0 = Inches(4.2)
    authors = [
        ("Selin Bardakcı",      "s.bardakci2022@gtu.edu.tr"),
        ("Bengisu Duru Göksu",  "bgoksu2022@gtu.edu.tr"),
        ("Zeynep Yiğit",        "z.yigit2022@gtu.edu.tr"),
    ]
    for i, (name, mail) in enumerate(authors):
        x = Inches(0.85 + i * 4.1)
        add_text(s, name, left=x, top=y0,
                 width=Inches(4), height=Inches(0.4),
                 size=16, bold=True, color=INK)
        add_text(s, mail, left=x, top=y0 + Inches(0.4),
                 width=Inches(4), height=Inches(0.3),
                 size=11, color=SUB)

    add_text(s, "Dept. of Computer Engineering · Gebze Technical University",
             left=Inches(0.85), top=Inches(5.5),
             width=Inches(12), height=Inches(0.4),
             size=14, color=SUB)
    add_text(s, "github.com/selin-bardakci/BigData-Github_Analysis",
             left=Inches(0.85), top=Inches(5.9),
             width=Inches(12), height=Inches(0.4),
             size=12, color=ACCENT, font="JetBrains Mono")
SLIDES.append(slide_title)


def slide_agenda():
    s = add_blank()
    header(s, "Agenda", kicker="Outline")
    items = [
        "Problem & motivation — what makes this hard",
        "Original approach (abandoned) — why it had to be redesigned",
        "Revised architecture — data sources, storage, compute",
        "Key design decisions & trade-offs",
        "Deliverables D1–D5 — preliminary results",
        "Interactive Web UI",
        "Conclusion & remaining work",
    ]
    add_bullets(s, items,
                left=Inches(0.85), top=Inches(1.8),
                width=Inches(12), height=Inches(5),
                size=20, line_spacing=1.6)
SLIDES.append(slide_agenda)


def slide_problem():
    s = add_blank()
    header(s, "Problem & Motivation", kicker="Why this matters")

    add_text(s, "Question",
             left=Inches(0.85), top=Inches(1.7),
             width=Inches(11), height=Inches(0.4),
             size=14, bold=True, color=ACCENT2)
    add_text(s,
             "Which programming technologies are emerging, which are declining, and "
             "what does the developer ecosystem behind them look like?",
             left=Inches(0.85), top=Inches(2.1),
             width=Inches(11.7), height=Inches(0.9),
             size=18, color=INK)

    add_text(s, "Scope",
             left=Inches(0.85), top=Inches(3.3),
             width=Inches(11), height=Inches(0.4),
             size=14, bold=True, color=ACCENT2)
    add_bullets(s, [
        "22.9M GitHub repositories — HuggingFace + GH Archive",
        "126 months of activity (January 2015 → June 2025)",
        "542 programming languages, 25,000 developers",
        "Five deliverables: trends, forecasts, migration, influence, communities",
    ], left=Inches(0.85), top=Inches(3.7),
       width=Inches(11.7), height=Inches(3),
       size=17, line_spacing=1.5)
SLIDES.append(slide_problem)


def slide_original_approach():
    s = add_blank()
    header(s, "Original Approach — and Why We Abandoned It",
           kicker="Phase 1 proposal")

    # left column: original
    add_text(s, "Phase-1 proposal", left=Inches(0.85), top=Inches(1.75),
             width=Inches(5.8), height=Inches(0.4),
             size=15, bold=True, color=BAD)
    add_bullets(s, [
        "Source: bigquery-public-data.github_repos (snapshot frozen ~2017)",
        "Modelling: ARIMA(2,1,2) for forecasting",
        "Co-authorship graph computed directly inside BigQuery SQL",
    ], left=Inches(0.85), top=Inches(2.15),
       width=Inches(5.7), height=Inches(2.4),
       size=14, bullet_color=BAD, line_spacing=1.45)

    # right column: problems
    add_text(s, "Two critical problems", left=Inches(7.0), top=Inches(1.75),
             width=Inches(5.8), height=Inches(0.4),
             size=15, bold=True, color=ACCENT4)
    add_bullets(s, [
        "Point-in-time labels — bytes column = state at single crawl date. "
        "Cross-joining onto every month produces a cross-sectional artefact, not a time-series.",
        "Coverage cut-off in 2017 — excludes the entire 2018–2025 rise of "
        "TypeScript, Rust, Go, Kotlin — the very languages most informative for forecasting.",
    ], left=Inches(7.0), top=Inches(2.15),
       width=Inches(5.7), height=Inches(3.3),
       size=14, bullet_color=ACCENT4, line_spacing=1.45)

    # bottom: outcome
    add_rect(s, Inches(0.85), Inches(5.6), Inches(11.7), Inches(1.1),
             fill=SURFACE, corner=True)
    add_text(s, "Decision",
             left=Inches(1.1), top=Inches(5.7),
             width=Inches(2), height=Inches(0.3),
             size=11, bold=True, color=ACCENT)
    add_text(s,
             "Replace both data sources, replace the forecasting model, "
             "and rebuild the pipeline on Spark.",
             left=Inches(1.1), top=Inches(6.0),
             width=Inches(11.3), height=Inches(0.6),
             size=16, color=INK)
SLIDES.append(slide_original_approach)


def slide_architecture():
    s = add_blank()
    header(s, "Revised Architecture", kicker="System overview")

    # 4-stage pipeline visual
    stages = [
        ("INGESTION", "HuggingFace 40M\n+ GH Archive\n→ GCS raw Parquet", ACCENT),
        ("PROCESSING", "6 Spark jobs on\nDataproc (us-east1)\nParquet + Snappy", ACCENT2),
        ("ANALYSIS",  "5 Jupyter notebooks\nread processed GCS\nfigures + tables", ACCENT3),
        ("DELIVERY",  "React + Recharts\nweb dashboard\nstatic JSON", ACCENT4),
    ]
    box_w = Inches(2.8)
    box_h = Inches(2.4)
    gap = Inches(0.3)
    total_w = box_w * 4 + gap * 3
    left0 = (prs.slide_width - total_w) // 2
    top = Inches(2.4)

    for i, (name, body, c) in enumerate(stages):
        x = left0 + i * (box_w + gap)
        add_rect(s, x, top, box_w, box_h, fill=SURFACE, corner=True)
        # accent bar at top
        add_rect(s, x, top, box_w, Inches(0.10), fill=c)
        add_text(s, f"0{i+1}", left=x, top=top + Inches(0.25),
                 width=box_w, height=Inches(0.4),
                 size=11, bold=True, color=c, align=PP_ALIGN.CENTER,
                 font="JetBrains Mono")
        add_text(s, name, left=x, top=top + Inches(0.6),
                 width=box_w, height=Inches(0.4),
                 size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_text(s, body, left=x + Inches(0.2),
                 top=top + Inches(1.15),
                 width=box_w - Inches(0.4), height=Inches(1.1),
                 size=11, color=SUB, align=PP_ALIGN.CENTER, line_spacing=1.35)
        if i < 3:
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + box_w + Inches(0.02),
                top + Inches(1.05),
                gap - Inches(0.04), Inches(0.4),
            )
            arrow.fill.solid(); arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()

    # caption
    add_text(s,
             "All inter-stage data is Apache Parquet + Snappy on GCS — "
             "columnar reads, predicate pushdown, and splittability for parallel Spark executors.",
             left=Inches(0.85), top=Inches(5.3),
             width=Inches(11.7), height=Inches(0.8),
             size=13, color=SUB, italic=True, line_spacing=1.4)

    # image placeholder for the full pipeline diagram from the report
    img_placeholder(s, Inches(0.85), Inches(6.1),
                    Inches(11.7), Inches(0.7),
                    "Fig. 1 — End-to-end pipeline diagram (from report)")
SLIDES.append(slide_architecture)


def slide_data_sources():
    s = add_blank()
    header(s, "Data Sources", kicker="What we ingest")

    # two side-by-side cards
    def card(left, color, title, body):
        add_rect(s, left, Inches(1.85), Inches(5.8), Inches(4.4),
                 fill=SURFACE, corner=True)
        add_rect(s, left, Inches(1.85), Inches(5.8), Inches(0.12),
                 fill=color)
        add_text(s, title, left=left + Inches(0.3), top=Inches(2.1),
                 width=Inches(5.4), height=Inches(0.45),
                 size=18, bold=True, color=INK)
        add_bullets(s, body, left=left + Inches(0.3), top=Inches(2.6),
                    width=Inches(5.3), height=Inches(3.6),
                    size=13, bullet_color=color, line_spacing=1.5)

    card(Inches(0.85), ACCENT,
         "HuggingFace github-repos-metadata-40M",
         [
            "22.9M repositories with Linguist language labels",
            "Creation dates spanning 8 years through 23 July 2025",
            "Sharded into 459 files → 459-way Spark parallelism",
            "Replaces the 2017-frozen BigQuery snapshot",
         ])
    card(Inches(6.85), ACCENT2,
         "GitHub Archive (githubarchive.month.*)",
         [
            "PushEvent counts per (repo, year_month) since 2015",
            "Live BigQuery dataset — updated monthly",
            "Cross-joined with HF labels on repo_name",
            "Replaces the frozen commits table from BigQuery",
         ])

    add_text(s,
             "Both sources were chosen to provide real time-series data — the "
             "fundamental defect of the original snapshot-based pipeline.",
             left=Inches(0.85), top=Inches(6.4),
             width=Inches(11.7), height=Inches(0.6),
             size=13, color=SUB, italic=True)
SLIDES.append(slide_data_sources)


def slide_spark_jobs():
    s = add_blank()
    header(s, "Spark Pipeline — Six Sequential Jobs", kicker="Processing")

    jobs = [
        ("d1_preprocess",          "Join Q1 labels with Q2 monthly pushes;\npre-birth filter to remove retroactive labels", ACCENT),
        ("d1_aggregate",           "Group by (language, year_month) →\nrepo_count, commit_count, total_bytes",                ACCENT),
        ("d1_features_cluster",    "5-D feature vectors → StandardScaler →\nKMeans (k=6, seed=42) on Spark MLlib",            ACCENT2),
        ("d2_forecasting",         "Holt-Winters ETS (damped trend + seasonal),\n24-month horizon for 30 top languages",      ACCENT3),
        ("d3_migration_graph",     "SQL pattern-matching on commit messages,\n70-tech allowlist, min edge weight = 3",       ACCENT4),
        ("d4d5_developer_graph",   "SHA-256 hash emails, build dev-dev graph,\nPageRank + Leiden community detection",        ACCENT2),
    ]
    cols, rows = 3, 2
    cell_w = Inches(4.0)
    cell_h = Inches(2.1)
    x0 = Inches(0.85); y0 = Inches(1.8); gap_x = Inches(0.15); gap_y = Inches(0.2)
    for i, (name, body, c) in enumerate(jobs):
        r, col = divmod(i, cols)
        x = x0 + col * (cell_w + gap_x)
        y = y0 + r * (cell_h + gap_y)
        add_rect(s, x, y, cell_w, cell_h, fill=SURFACE, corner=True)
        add_rect(s, x, y, Inches(0.08), cell_h, fill=c)
        add_text(s, name, left=x + Inches(0.25), top=y + Inches(0.15),
                 width=cell_w - Inches(0.3), height=Inches(0.4),
                 size=14, bold=True, color=INK, font="JetBrains Mono")
        add_text(s, body, left=x + Inches(0.25), top=y + Inches(0.65),
                 width=cell_w - Inches(0.3), height=Inches(1.4),
                 size=12, color=SUB, line_spacing=1.35)

    add_text(s,
             "Job orchestration: sequential + wait-poll via gcloud dataproc jobs wait — "
             "parallel --async caused corrupt Parquet footer errors.",
             left=Inches(0.85), top=Inches(6.55),
             width=Inches(11.7), height=Inches(0.5),
             size=12, color=ACCENT4, italic=True)
SLIDES.append(slide_spark_jobs)


def slide_tradeoffs_1():
    s = add_blank()
    header(s, "Design Decisions (1/3) — Data & Forecasting",
           kicker="What we tried and rejected")

    rows = [
        ("Language data source", "HuggingFace 40M (2025)",
         "BigQuery snapshot (2017)",
         "BQ is point-in-time; HF provides timestamps + 8 more years"),
        ("Commit data source", "GH Archive PushEvents",
         "BigQuery frozen commits",
         "BQ commits table frozen at 2017; GH Archive is live and monthly"),
        ("Forecasting model", "Holt-Winters ETS (damped)",
         "ARIMA(2,1,2) · Prophet",
         "ARIMA: flat lines for multi-year growth · Prophet: CmdStan won't build on Dataproc"),
        ("Storage format", "Parquet + Snappy",
         "CSV, ORC, AVRO",
         "Columnar reads + predicate pushdown; Snappy fast for compute-bound jobs"),
    ]
    _decision_table(s, rows)
SLIDES.append(slide_tradeoffs_1)


def slide_tradeoffs_2():
    s = add_blank()
    header(s, "Design Decisions (2/3) — Graph Analytics",
           kicker="What we tried and rejected")

    rows = [
        ("Community detection", "Leiden",
         "Louvain",
         "Leiden guarantees well-connected communities; Louvain can split sub-clusters"),
        ("Influence metric", "PageRank + Betweenness",
         "PageRank + HITS",
         "HITS hub_score ≡ authority_score on undirected graphs (mathematically degenerate)"),
        ("Developer selection", "Distinct repository count",
         "Total commit count",
         "Total commits selects CI bots (e.g. 789K commits); distinct repos finds real contributors"),
        ("Megarepo exclusion", "MAX_DEVS_PER_REPO = 200",
         "No exclusion",
         "Linux/TensorFlow contributor lists cause O(N²) join → 50M+ edge rows from one repo"),
    ]
    _decision_table(s, rows)
SLIDES.append(slide_tradeoffs_2)


def slide_tradeoffs_3():
    s = add_blank()
    header(s, "Design Decisions (3/3) — Infrastructure",
           kicker="What we tried and rejected")

    rows = [
        ("Job orchestration", "Sequential + wait-poll",
         "Parallel --async submission",
         "Parallel jobs corrupted GCS outputs — downstream read paths still being written"),
        ("Driver-side libraries", "_ensure_module() + pip install at runtime",
         "Custom Dataproc image",
         "Avoids image rebuild while keeping Spark scripts self-contained"),
        ("Dataproc zone", "us-east1 (auto-zone)",
         "us-central1-a · us-central1-f",
         "Both us-central1 zones returned 'Out of capacity' twice consecutively"),
        ("Data cutoff", "year_month ≤ 2025-06",
         "2025-12 or no cap",
         "HF dataset ends 23 July 2025; post-cutoff months show 49% artificial repo-count drop"),
    ]
    _decision_table(s, rows)
SLIDES.append(slide_tradeoffs_3)


def _decision_table(slide, rows):
    """Render a 4-column decision table: point | chosen | rejected | reason."""
    headers = ("Decision point", "Chosen ✓", "Rejected ✗", "Reason")
    col_w = [Inches(2.6), Inches(2.7), Inches(2.7), Inches(4.4)]
    left = Inches(0.85)
    top = Inches(1.8)
    row_h = Inches(1.05)

    # header row
    x = left
    for i, h in enumerate(headers):
        add_rect(slide, x, top, col_w[i], Inches(0.45),
                 fill=RGBColor(0x1F, 0x27, 0x33))
        color = ACCENT if i == 1 else BAD if i == 2 else INK
        add_text(slide, h, left=x + Inches(0.15), top=top + Inches(0.08),
                 width=col_w[i] - Inches(0.2), height=Inches(0.3),
                 size=12, bold=True, color=color)
        x += col_w[i]

    # data rows
    y = top + Inches(0.5)
    for r, (point, chosen, rejected, reason) in enumerate(rows):
        # alternate row background
        if r % 2 == 0:
            add_rect(slide, left, y, sum(col_w, Inches(0)), row_h,
                     fill=SURFACE)
        x = left
        cells = [(point, INK, True), (chosen, GOOD, True),
                 (rejected, BAD, False), (reason, SUB, False)]
        for i, (text, color, bold) in enumerate(cells):
            font = "JetBrains Mono" if i in (1, 2) else "Inter"
            add_text(slide, text,
                     left=x + Inches(0.15), top=y + Inches(0.12),
                     width=col_w[i] - Inches(0.2), height=row_h - Inches(0.2),
                     size=11, bold=bold, color=color, font=font,
                     line_spacing=1.25)
            x += col_w[i]
        y += row_h


def slide_d1():
    s = add_blank()
    header(s, "D1 — Language Growth Trends", kicker="Deliverable 1 results")

    add_bullets(s, [
        "53,630 monthly rows across 542 languages over 126 months",
        "Top-10 by mean monthly repo count: JavaScript (~133K), Python (~96K), "
        "TypeScript (~47K), Java (~72K), PHP (~38K)",
        "K-Means (k=6) partitions into trajectory archetypes:",
        "   • Fast-growing modern — TypeScript, Rust, Go, Kotlin, Dart",
        "   • Dominant incumbents — JavaScript, Python, Java",
        "   • Stable systems — C, C++, Objective-C",
        "   • Declining legacy — Ruby, PHP, Perl, CoffeeScript",
        "   • Niche/academic — R, Haskell, Julia, Erlang, OCaml",
        "   • Long-tail — ~430 low-activity languages",
    ], left=Inches(0.85), top=Inches(1.85),
       width=Inches(7), height=Inches(5),
       size=13, line_spacing=1.5)

    img_placeholder(s, Inches(8.2), Inches(1.85),
                    Inches(4.4), Inches(2.6),
                    "outputs/d1_timeseries.png")
    img_placeholder(s, Inches(8.2), Inches(4.6),
                    Inches(4.4), Inches(2.3),
                    "outputs/d1_cluster_pca.png")
SLIDES.append(slide_d1)


def slide_d2():
    s = add_blank()
    header(s, "D2 — 24-Month Holt-Winters Forecasts",
           kicker="Deliverable 2 results")

    add_text(s, "Projected 24-month change from June 2025 baseline",
             left=Inches(0.85), top=Inches(1.8),
             width=Inches(11), height=Inches(0.4),
             size=13, bold=True, color=ACCENT2)

    rows = [
        ("Nix",           "+27%", GOOD,    "Reproducible build adoption"),
        ("C#",            "+26%", GOOD,    "Unity + .NET MAUI segments"),
        ("Go, Rust, Julia","+18%", GOOD,    "Systems & ML momentum"),
        ("TypeScript, Python", "+10%", GOOD, "Sustained expansion at scale"),
        ("JavaScript",    "+4%",  ACCENT4, "Plateau — market saturation"),
        ("Vim Script, SCSS, Objective-C, Ruby, Lua, PHP",
         "−6 to −27%", BAD, "Consistent multi-year decline"),
    ]
    y = Inches(2.4)
    for lang, delta, color, note in rows:
        add_rect(s, Inches(0.85), y, Inches(7), Inches(0.5),
                 fill=SURFACE, corner=True)
        add_text(s, lang, left=Inches(1.05), top=y + Inches(0.1),
                 width=Inches(4.0), height=Inches(0.4),
                 size=12, color=INK, bold=True)
        add_text(s, delta, left=Inches(4.8), top=y + Inches(0.1),
                 width=Inches(1.4), height=Inches(0.4),
                 size=12, color=color, bold=True, font="JetBrains Mono",
                 align=PP_ALIGN.RIGHT)
        add_text(s, note, left=Inches(6.4), top=y + Inches(0.1),
                 width=Inches(1.4), height=Inches(0.4),
                 size=10, color=SUB)
        y += Inches(0.6)

    img_placeholder(s, Inches(8.2), Inches(1.85),
                    Inches(4.4), Inches(2.5),
                    "outputs/d2_top6_forecasts.png")
    img_placeholder(s, Inches(8.2), Inches(4.5),
                    Inches(4.4), Inches(2.5),
                    "outputs/d2_forecast_summary.png")
SLIDES.append(slide_d2)


def slide_d3():
    s = add_blank()
    header(s, "D3 — Technology Migration Graph",
           kicker="Deliverable 3 results")

    add_bullets(s, [
        "Pattern-matched commit messages: \"migrated from X to Y\", "
        "\"switching from X to Y\", \"moving from X to Y\"",
        "70-technology allowlist + minimum edge weight = 3 → "
        "16 nodes, 11 high-confidence edges",
        "Top migration attractors: CircleCI (n=24), Jest (n=18), Go (n=12)",
        "Strongest sources: Mocha, Travis CI — replaced by modern CI/CD and test runners",
        "Edge coverage limitation: GH Archive commit messages are free-form. "
        "Phase-3 plan: extend to PR/issue text + CHANGELOG parsing, possibly NLP-based.",
    ], left=Inches(0.85), top=Inches(1.85),
       width=Inches(7.2), height=Inches(5),
       size=14, line_spacing=1.55)

    img_placeholder(s, Inches(8.4), Inches(1.85),
                    Inches(4.2), Inches(4.8),
                    "outputs/d3_network.png   (or outputs/d3_sankey.html)")
SLIDES.append(slide_d3)


def slide_d4():
    s = add_blank()
    header(s, "D4 — Developer Influence (PageRank)",
           kicker="Deliverable 4 results")

    add_bullets(s, [
        "25,000 developers · mean degree 7.7 · max degree 211",
        "9,625 isolated nodes (38.5%) — even after bot filtering, many developers "
        "contribute to a single project only",
        "PageRank (α = 0.85) weighted by shared repositories — identifies the "
        "connectors whose removal would most degrade reachability",
        "Betweenness centrality complements PageRank — finds bridge nodes between disjoint clusters",
        "Weak correlation between PageRank and commit count (r = 0.277) — "
        "high-commit bots have near-zero PageRank; well-connected cross-project contributors rank highly",
        "Privacy: emails irreversibly SHA-256 hashed at DataFrame creation; no raw email ever written",
    ], left=Inches(0.85), top=Inches(1.85),
       width=Inches(7.2), height=Inches(5),
       size=13, line_spacing=1.5)

    img_placeholder(s, Inches(8.4), Inches(1.85),
                    Inches(4.2), Inches(2.4),
                    "outputs/d4_top20_pagerank.png")
    img_placeholder(s, Inches(8.4), Inches(4.4),
                    Inches(4.2), Inches(2.4),
                    "outputs/d4_pagerank_vs_commits.png")
SLIDES.append(slide_d4)


def slide_d5():
    s = add_blank()
    header(s, "D5 — Developer Community Detection",
           kicker="Deliverable 5 results")

    add_bullets(s, [
        "Leiden algorithm yielded 10,255 communities from 25,000 developers",
        "Largest community: 1,679 members — dense web-ecosystem cluster "
        "(JavaScript / TypeScript repositories)",
        "Median community size: 1 member (singleton) — 93.9% of all communities",
        "Heavy-tailed distribution matches open-source collaboration theory: "
        "few large maintained projects, many small isolated efforts",
        "Top communities labelled by dominant GitHub organisations: "
        "fastify · conda-forge · learn-co-students · kubernetes · google · facebook",
        "Leiden chosen over Louvain to guarantee well-connectedness (no internally-disconnected partitions) "
        "and reproducibility (seed=42)",
    ], left=Inches(0.85), top=Inches(1.85),
       width=Inches(7.2), height=Inches(5),
       size=13, line_spacing=1.5)

    img_placeholder(s, Inches(8.4), Inches(1.85),
                    Inches(4.2), Inches(2.4),
                    "outputs/d5_community_sizes.png")
    img_placeholder(s, Inches(8.4), Inches(4.4),
                    Inches(4.2), Inches(2.4),
                    "outputs/d5_size_vs_influence.png")
SLIDES.append(slide_d5)


def slide_webui():
    s = add_blank()
    header(s, "Interactive Web UI", kicker="Beyond the report")

    add_bullets(s, [
        "React + Vite + Tailwind dashboard built on top of the Spark outputs",
        "Spark runs once offline; UI reads slim JSON (~840 KB total) — no backend, no Spark round-trip",
        "Client-side filters: language chips, year-range slider, log/linear, normalised growth index",
        "Pre-computed time windows for the migration graph (instant switching, no recomputation)",
        "Dark-mode design, 6 pages — Overview + D1-D5",
    ], left=Inches(0.85), top=Inches(1.85),
       width=Inches(11.7), height=Inches(2.4),
       size=14, line_spacing=1.5)

    # screenshot placeholders — three side by side
    img_placeholder(s, Inches(0.85), Inches(4.4),
                    Inches(3.9), Inches(2.5),
                    "screenshots/ui_overview.png")
    img_placeholder(s, Inches(4.95), Inches(4.4),
                    Inches(3.9), Inches(2.5),
                    "screenshots/ui_d1_trends.png")
    img_placeholder(s, Inches(9.05), Inches(4.4),
                    Inches(3.6), Inches(2.5),
                    "screenshots/ui_d3_sankey.png")
SLIDES.append(slide_webui)


def slide_conclusion():
    s = add_blank()
    header(s, "Conclusion & Remaining Work", kicker="Next steps")

    add_text(s, "What works today", left=Inches(0.85), top=Inches(1.8),
             width=Inches(11), height=Inches(0.4),
             size=15, bold=True, color=GOOD)
    add_bullets(s, [
        "End-to-end pipeline runs successfully on Dataproc (us-east1)",
        "All five deliverables produce quantitative outputs",
        "Main findings: TypeScript / Rust lead growth; CircleCI + Jest are migration attractors; "
        "developer collaboration follows a heavy-tailed community structure",
    ], left=Inches(0.85), top=Inches(2.2),
       width=Inches(11.7), height=Inches(2),
       size=13, bullet_color=GOOD, line_spacing=1.5)

    add_text(s, "Phase-3 plan", left=Inches(0.85), top=Inches(4.2),
             width=Inches(11), height=Inches(0.4),
             size=15, bold=True, color=ACCENT)
    add_bullets(s, [
        "Enrich D3 migration graph with PR / issue text and CHANGELOG parsing",
        "Add Holt-Winters confidence intervals to the D2 notebook",
        "Cross-join D4 developer PageRank with D1 language clusters",
        "Evaluate alignment of D5 communities with D1 language clusters",
        "Polish the interactive Web UI and wire it to the live Spark outputs",
    ], left=Inches(0.85), top=Inches(4.6),
       width=Inches(11.7), height=Inches(2.4),
       size=13, line_spacing=1.5)
SLIDES.append(slide_conclusion)


def slide_thanks():
    s = add_blank()
    add_rect(s, 0, 0, prs.slide_width, Inches(0.06), fill=ACCENT)
    add_rect(s, 0, Inches(0.06), prs.slide_width, Inches(0.04), fill=ACCENT2)

    add_text(s, "Thank you.",
             left=Inches(0.85), top=Inches(2.6),
             width=Inches(12), height=Inches(1.2),
             size=64, bold=True, color=INK)
    add_text(s, "Questions, comments, suggestions welcome.",
             left=Inches(0.85), top=Inches(3.8),
             width=Inches(12), height=Inches(0.6),
             size=20, color=ACCENT2, italic=True)

    add_text(s, "github.com/selin-bardakci/BigData-Github_Analysis",
             left=Inches(0.85), top=Inches(5.0),
             width=Inches(12), height=Inches(0.4),
             size=14, color=ACCENT, font="JetBrains Mono")
    add_text(s, "Selin Bardakcı · Bengisu Duru Göksu · Zeynep Yiğit "
                "· Gebze Technical University · CSE",
             left=Inches(0.85), top=Inches(5.45),
             width=Inches(12), height=Inches(0.4),
             size=12, color=SUB)
SLIDES.append(slide_thanks)


# ── Build all slides ───────────────────────────────────────────────────────
total = len(SLIDES)
for i, builder in enumerate(SLIDES, start=1):
    before = len(prs.slides)
    builder()
    # footer on every slide except title and thanks
    if 1 < i < total:
        footer(prs.slides[before], i, total)

prs.save(OUT)
print(f"Wrote {total} slides → {OUT}")
